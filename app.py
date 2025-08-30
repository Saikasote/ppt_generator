import io, json, os, re, textwrap, tempfile, uuid
from typing import List, Dict, Optional

import httpx
from fastapi import FastAPI, UploadFile, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
from pptx.util import Inches, Pt
from starlette.background import BackgroundTask

APP_DIR = os.path.dirname(os.path.abspath(__file__))
FRONTEND_DIR = os.path.join(APP_DIR, "frontend")

app = FastAPI(title="Auto PPT Generator", version="1.0.0")

# Serve frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # tighten in prod
    allow_methods=["*"],
    allow_headers=["*"],
)
app.mount("/static", StaticFiles(directory=FRONTEND_DIR), name="static")


@app.get("/", response_class=HTMLResponse)
async def root():
    with open(os.path.join(FRONTEND_DIR, "index.html"), "r", encoding="utf-8") as f:
        return f.read()


# ---------------------------
# Utilities
# ---------------------------

def clean_filename(name: str) -> str:
    name = name.strip() or "presentation"
    return re.sub(r"[^a-zA-Z0-9_\-]", "_", name)

def wrap_lines(text: str, width: int = 80) -> str:
    return "\n".join(textwrap.wrap(text, width=width))

def heuristic_outline(text: str, guidance: str = "") -> List[Dict]:
    """
    Split long text into slide-friendly sections without LLM.
    - Detect headings by blank lines / markdown style
    - Otherwise chunk into ~7-10 bullet items per slide
    """
    lines = [ln.strip() for ln in text.splitlines()]
    # Collapse multiple blanks
    blocks, buf = [], []
    for ln in lines:
        if not ln:
            if buf:
                blocks.append("\n".join(buf))
                buf = []
        else:
            buf.append(ln)
    if buf:
        blocks.append("\n".join(buf))

    slides = []
    for i, block in enumerate(blocks, start=1):
        # If looks like a heading (markdown '# ' or Title line + following text)
        title = None
        body = block
        md_head = re.match(r"^(#+)\s+(.*)", block.split("\n", 1)[0])
        if md_head:
            title = md_head.group(2).strip()
            body = "\n".join(block.split("\n")[1:])
        else:
            # First line as title if short
            first = block.split("\n", 1)[0].strip()
            if 0 < len(first) <= 70:
                title = first
                body = "\n".join(block.split("\n")[1:])

        bullets = []
        for ln in body.splitlines():
            ln = ln.strip()
            if not ln:
                continue
            # strip markdown bullets
            ln = re.sub(r"^(\-|\*|\d+\.)\s+", "", ln)
            bullets.append(ln)

        # chunk bullets into multiple slides if too long
        if not bullets:
            bullets = [body] if body else []

        chunk_size = 8
        for c in range(0, len(bullets), chunk_size):
            part = bullets[c:c+chunk_size]
            slides.append({
                "title": (title or f"Slide {len(slides)+1}") + (f" ({c//chunk_size+1})" if len(bullets) > chunk_size else ""),
                "bullets": part,
                "notes": guidance or ""
            })

    if not slides:
        slides = [{"title": "Overview", "bullets": [text], "notes": guidance or ""}]
    return slides


async def llm_outline(provider: str, api_key: str, text: str, guidance: str, model: Optional[str]) -> Optional[List[Dict]]:
    """
    Ask an LLM to return a JSON outline:
    [{"title": "...", "bullets": ["..."], "notes": "..."}]
    Supports openai | anthropic | gemini. Returns None on failure.
    """
    provider = (provider or "").lower()
    model = (model or "").strip()
    sys_prompt = (
        "You convert long text into a slide outline. "
        "Return STRICT JSON only (no prose), as an array of objects: "
        '{"title": str, "bullets": [str, ...], "notes": str}. '
        "Aim for 6-15 slides. Keep bullets concise (<= 14 words each)."
    )
    user_prompt = f"Guidance: {guidance or 'none'}\n\nContent:\n{text}"

    try:
        if provider == "openai":
            url = "https://api.openai.com/v1/chat/completions"
            use_model = model or "gpt-4o-mini"
            headers = {"Authorization": f"Bearer {api_key}"}
            payload = {
                "model": use_model,
                "temperature": 0.3,
                "response_format": {"type": "json_object"},
                "messages": [
                    {"role": "system", "content": sys_prompt},
                    {"role": "user", "content": user_prompt}
                ],
            }
            async with httpx.AsyncClient(timeout=60) as client:
                r = await client.post(url, headers=headers, json=payload)
                r.raise_for_status()
                content = r.json()["choices"][0]["message"]["content"]
                data = json.loads(content)
                # accept {"slides":[...]} or [...]
                slides = data.get("slides") if isinstance(data, dict) else data
                if isinstance(slides, list):
                    return slides

        elif provider == "anthropic":
            url = "https://api.anthropic.com/v1/messages"
            use_model = model or "claude-3-5-sonnet-latest"
            headers = {
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json"
            }
            payload = {
                "model": use_model,
                "max_tokens": 2000,
                "temperature": 0.3,
                "system": sys_prompt,
                "messages": [{"role": "user", "content": user_prompt}],
            }
            async with httpx.AsyncClient(timeout=60) as client:
                r = await client.post(url, headers=headers, json=payload)
                r.raise_for_status()
                msg = r.json()["content"][0]["text"]
                # try parse code block json
                m = re.search(r"\{.*\}|\[.*\]", msg, re.S)
                content = m.group(0) if m else msg
                data = json.loads(content)
                slides = data.get("slides") if isinstance(data, dict) else data
                if isinstance(slides, list):
                    return slides

        elif provider == "gemini":
            # Gemini uses key as query param
            use_model = model or "gemini-1.5-flash"
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{use_model}:generateContent?key={api_key}"
            payload = {
                "contents": [{
                    "parts": [{"text": f"{sys_prompt}\n\n{user_prompt}"}]
                }],
                "generationConfig": {"temperature": 0.3}
            }
            async with httpx.AsyncClient(timeout=60) as client:
                r = await client.post(url, json=payload)
                r.raise_for_status()
                candidates = r.json().get("candidates", [])
                text_out = candidates[0]["content"]["parts"][0]["text"] if candidates else "[]"
                m = re.search(r"\{.*\}|\[.*\]", text_out, re.S)
                content = m.group(0) if m else text_out
                data = json.loads(content)
                slides = data.get("slides") if isinstance(data, dict) else data
                if isinstance(slides, list):
                    return slides

    except Exception:
        # swallow and fallback
        return None

    return None


def extract_template_images(prs: Presentation) -> List[bytes]:
    """Collect image blobs from the template (dedup by size)."""
    seen = set()
    imgs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                blob = shape.image.blob
                key = (len(blob), hash(blob[:64]))
                if key not in seen:
                    seen.add(key)
                    imgs.append(blob)
    return imgs


def add_picture_if_placeholder(slide, pic_blob: Optional[bytes]):
    """
    If the layout has a picture placeholder, put the template image into it.
    Otherwise, add a small decorative image in bottom-right.
    """
    if not pic_blob:
        return

    pic_placeholder = None
    for shp in slide.shapes:
        if shp.is_placeholder and getattr(shp.placeholder_format, "type", None) == 18:  # PICTURE
            pic_placeholder = shp
            break

    if pic_placeholder:
        # Replace placeholder
        try:
            pic_placeholder.insert_picture(io.BytesIO(pic_blob))
            return
        except Exception:
            pass

    # Otherwise, add as floating image
    left = slide.slide_width - Inches(2.0) if hasattr(slide, "slide_width") else Inches(8.0)
    # Fallback positions:
    try:
        slide.shapes.add_picture(io.BytesIO(pic_blob), Inches(8.0), Inches(5.0), width=Inches(1.8))
    except Exception:
        # ignore failures silently
        pass


def build_presentation(template_path: Optional[str], slides_data: List[Dict], out_path: str):
    prs = Presentation(template_path) if template_path else Presentation()
    # prefer title+content layout if exists
    layout = None
    for idx in range(len(prs.slide_layouts)):
        l = prs.slide_layouts[idx]
        # try to find layout with title + body placeholder
        if any(sh.is_placeholder and sh.placeholder_format.type == 1 for sh in l.shapes) and \
           any(sh.is_placeholder and sh.placeholder_format.type in (2, 7) for sh in l.shapes):
            layout = l
            break
    if not layout:
        layout = prs.slide_layouts[0]

    template_images = extract_template_images(prs)
    template_image = template_images[0] if template_images else None

    for s in slides_data:
        slide = prs.slides.add_slide(layout)
        # Title
        try:
            slide.shapes.title.text = s.get("title") or "Slide"
        except Exception:
            pass

        # Body / bullets
        body = None
        for shp in slide.shapes:
            if shp.is_placeholder and getattr(shp.placeholder_format, "type", None) in (2, 7):  # BODY or CONTENT
                body = shp
                break
        if body:
            tf = body.text_frame
            tf.clear()
            bullets = s.get("bullets") or []
            if not bullets and s.get("title"):
                bullets = [s["title"]]
            for i, b in enumerate(bullets):
                if i == 0:
                    tf.text = b
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0
            # Try to keep text readable
            for p in tf.paragraphs:
                for run in p.runs:
                    try:
                        run.font.size = Pt(20)
                    except Exception:
                        pass

        # Notes
        notes = s.get("notes") or ""
        try:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = wrap_lines(notes, 100)
        except Exception:
            pass

        # Images from template
        add_picture_if_placeholder(slide, template_image)

    prs.save(out_path)


# ---------------------------
# API
# ---------------------------

@app.post("/generate-ppt")
async def generate_ppt(
    text_content: str = Form(...),
    guidance: str = Form(""),
    llm_provider: str = Form(""),
    api_key: str = Form(""),
    model: str = Form(""),
    filename: str = Form("presentation"),
    template_file: UploadFile | None = None,
):
    """
    Core endpoint:
    - Uses LLM (if provider+key) to build a slide outline JSON.
    - Otherwise uses heuristic splitting.
    - Applies the uploaded template styles to slides.
    - Returns a PPTX.
    """
    # Basic limits
    if len(text_content) > 100_000:
        return HTMLResponse("Text too long (limit ~100k chars).", status_code=413)

    # Never store API keys
    provider = (llm_provider or "").strip()
    key = (api_key or "").strip()

    # Save template to temp file if provided
    template_path = None
    if template_file:
        if not template_file.filename.lower().endswith((".pptx", ".potx")):
            return HTMLResponse("Template must be .pptx or .potx", status_code=400)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(template_file.filename)[1])
        tmp.write(await template_file.read())
        tmp.close()
        template_path = tmp.name

    # Get outline: LLM first, else heuristic
    slides_data = None
    if provider and key:
        slides_data = await llm_outline(provider, key, text_content, guidance, model)
    if not slides_data:
        slides_data = heuristic_outline(text_content, guidance)

    # Build PPT
    outname = clean_filename(filename) + ".pptx"
    outpath = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4().hex}-{outname}")
    build_presentation(template_path, slides_data, outpath)

    # cleanup template after response is sent
    def cleanup():
        try:
            if template_path and os.path.exists(template_path):
                os.remove(template_path)
            if os.path.exists(outpath):
                os.remove(outpath)
        except Exception:
            pass

    return FileResponse(
        outpath,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=outname,
        background=BackgroundTask(cleanup),
    )
